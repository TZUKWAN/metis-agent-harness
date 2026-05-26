"""Document generation tools: flowchart, Excel, Word, PPT."""

from __future__ import annotations

from pathlib import Path

from metis.security.paths import is_write_denied
from metis.tools.registry import ToolRegistry
from metis.tools.spec import ToolContext, ToolPermissionLevel, ToolSpec


def register_document_tools(registry: ToolRegistry, *, workspace: str = ".") -> None:
    root = Path(workspace).resolve()

    def create_flowchart(args: dict, context: ToolContext) -> dict:
        """Create a Mermaid flowchart markdown file from node/edge definitions."""
        path = args["path"]
        direction = args.get("direction", "TD")
        nodes = args.get("nodes", [])
        edges = args.get("edges", [])
        title = args.get("title", "Flowchart")

        lines = [f"## {title}", "", f"```mermaid", f"flowchart {direction}"]
        for node in nodes:
            node_id = node.get("id", "")
            label = node.get("label", node_id)
            shape = node.get("shape", "rect")
            if shape == "rounded":
                lines.append(f"    {node_id}({label})")
            elif shape == "diamond":
                lines.append(f"    {node_id}{{{label}}}")
            elif shape == "circle":
                lines.append(f"    {node_id}(({label}))")
            else:
                lines.append(f"    {node_id}[{label}]")
        for edge in edges:
            src = edge.get("from", "")
            dst = edge.get("to", "")
            label = edge.get("label", "")
            if label:
                lines.append(f"    {src} -->|{label}| {dst}")
            else:
                lines.append(f"    {src} --> {dst}")
        lines.append("```")

        content = "\n".join(lines)
        full_path = root / path
        if is_write_denied(str(full_path)):
            return {"error": f"Write denied for security: {path}"}
        full_path.parent.mkdir(parents=True, exist_ok=True)
        full_path.write_text(content, encoding="utf-8")
        return {"path": str(full_path), "size": len(content), "format": "mermaid"}

    def create_spreadsheet(args: dict, context: ToolContext) -> dict:
        """Create an Excel spreadsheet from headers and rows."""
        from openpyxl import Workbook

        path = args["path"]
        headers = args.get("headers", [])
        rows = args.get("rows", [])
        sheet_name = args.get("sheet_name", "Sheet1")

        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name
        if headers:
            ws.append(headers)
        for row in rows:
            ws.append(row)

        full_path = root / path
        if is_write_denied(str(full_path)):
            return {"error": f"Write denied for security: {path}"}
        full_path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(full_path))
        return {"path": str(full_path), "rows": len(rows), "columns": len(headers), "format": "xlsx"}

    def create_document(args: dict, context: ToolContext) -> dict:
        """Create a Word document with title, paragraphs, and optional headings."""
        from docx import Document

        path = args["path"]
        title = args.get("title", "")
        sections = args.get("sections", [])

        doc = Document()
        if title:
            doc.add_heading(title, level=0)
        for section in sections:
            heading = section.get("heading", "")
            body = section.get("body", "")
            level = section.get("level", 1)
            if heading:
                doc.add_heading(heading, level=level)
            if body:
                doc.add_paragraph(body)

        full_path = root / path
        if is_write_denied(str(full_path)):
            return {"error": f"Write denied for security: {path}"}
        full_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(full_path))
        return {"path": str(full_path), "sections": len(sections), "format": "docx"}

    def create_presentation(args: dict, context: ToolContext) -> dict:
        """Create a PowerPoint presentation with slides."""
        from pptx import Presentation

        path = args["path"]
        slides = args.get("slides", [])

        prs = Presentation()
        for slide_data in slides:
            title = slide_data.get("title", "")
            body = slide_data.get("body", "")
            layout_idx = min(int(slide_data.get("layout", 0)), len(prs.slide_layouts) - 1)
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            if title and slide.shapes.title:
                slide.shapes.title.text = title
            if body:
                placeholders = slide.placeholders
                body_placeholder = None
                for ph in placeholders:
                    if ph.placeholder_format.idx == 1:
                        body_placeholder = ph
                        break
                if body_placeholder:
                    body_placeholder.text = body

        full_path = root / path
        if is_write_denied(str(full_path)):
            return {"error": f"Write denied for security: {path}"}
        full_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(full_path))
        return {"path": str(full_path), "slides": len(slides), "format": "pptx"}

    registry.register(
        ToolSpec(
            name="create_flowchart",
            description="Create a flowchart as a Mermaid markdown file. Provide nodes (id, label, shape) and edges (from, to, label). Shapes: rect, rounded, diamond, circle.",
            parameters={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "minLength": 1},
                    "title": {"type": "string"},
                    "direction": {"type": "string", "enum": ["TD", "LR", "BT", "RL"]},
                    "nodes": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "id": {"type": "string"},
                                "label": {"type": "string"},
                                "shape": {"type": "string", "enum": ["rect", "rounded", "diamond", "circle"]},
                            },
                            "required": ["id", "label"],
                        },
                    },
                    "edges": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "from": {"type": "string"},
                                "to": {"type": "string"},
                                "label": {"type": "string"},
                            },
                            "required": ["from", "to"],
                        },
                    },
                },
                "required": ["path"],
                "additionalProperties": False,
            },
            handler=create_flowchart,
            category="documents",
            side_effect="write",
            permission_level=ToolPermissionLevel.WORKSPACE_WRITE.value,
        )
    )
    registry.register(
        ToolSpec(
            name="create_spreadsheet",
            description="Create an Excel file (.xlsx) with headers and rows. Example: headers=['Name','Age'], rows=[['Alice',30],['Bob',25]].",
            parameters={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "minLength": 1},
                    "sheet_name": {"type": "string"},
                    "headers": {"type": "array", "items": {"type": "string"}},
                    "rows": {"type": "array", "items": {"type": "array"}},
                },
                "required": ["path", "headers", "rows"],
                "additionalProperties": False,
            },
            handler=create_spreadsheet,
            category="documents",
            side_effect="write",
            permission_level=ToolPermissionLevel.WORKSPACE_WRITE.value,
        )
    )
    registry.register(
        ToolSpec(
            name="create_document",
            description="Create a Word document (.docx) with title and sections. Each section has a heading and body text.",
            parameters={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "minLength": 1},
                    "title": {"type": "string"},
                    "sections": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {"type": "string"},
                                "body": {"type": "string"},
                                "level": {"type": "integer", "minimum": 1, "maximum": 4},
                            },
                        },
                    },
                },
                "required": ["path"],
                "additionalProperties": False,
            },
            handler=create_document,
            category="documents",
            side_effect="write",
            permission_level=ToolPermissionLevel.WORKSPACE_WRITE.value,
        )
    )
    registry.register(
        ToolSpec(
            name="create_presentation",
            description="Create a PowerPoint file (.pptx) with slides. Each slide has a title and body text.",
            parameters={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "minLength": 1},
                    "slides": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "body": {"type": "string"},
                                "layout": {"type": "integer"},
                            },
                        },
                    },
                },
                "required": ["path", "slides"],
                "additionalProperties": False,
            },
            handler=create_presentation,
            category="documents",
            side_effect="write",
            permission_level=ToolPermissionLevel.WORKSPACE_WRITE.value,
        )
    )
